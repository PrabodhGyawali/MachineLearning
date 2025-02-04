from pathlib import Path
from pptx import Presentation # Must do a pip install for this
import sys
import threading
from queue import Queue
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
import logging

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def find_pptx_files(directory_path):
    """Find all .pptx files in the given directory"""
    try:
        dir_path = Path(directory_path)

        if not dir_path.exists():
            raise FileNotFoundError(f"Directory not found {directory_path}")

        pptx_files = list(dir_path.glob("*.pptx"))
        pptx_files.sort()

        if not pptx_files:
            logger.warning(f"No .pptx files found in {directory_path}")
            return []

        logger.info(f"Found {len(pptx_files)} PowerPoint files")
        return pptx_files

    except Exception as e:
        logger.error(f"Error processing directory: {str(e)}")
        return []

def extract_text_from_pptx(pptx_path):
    """Extract text from a single PowerPoint file"""
    try:
        prs = Presentation(pptx_path)
        text_content = []

        # Add file name as header
        text_content.append(f"=== Content from: {pptx_path.name} ===\n")

        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n--- Slide {slide_num} ---\n")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = " ".join(run.text.strip() for run in paragraph.runs)
                    if paragraph_text:
                        text_content.append(paragraph_text)

        # Create output file path
        output_path = pptx_path.with_suffix('.txt')

        # Write content to file
        with output_path.open('w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))

        logger.info(f"Successfully processed {pptx_path.name}")
        return True

    except Exception as e:
        logger.error(f"Error processing {pptx_path.name}: {str(e)}")
        return False

def process_directory(directory_path, max_workers=4):
    """Process all PowerPoint files in parallel"""
    pptx_files = find_pptx_files(directory_path)
    if not pptx_files:
        return

    start_time = time.time()
    successful = 0
    failed = 0

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Submit all tasks and create a future-to-filename mapping
        future_to_file = {
            executor.submit(extract_text_from_pptx, pptx_file): pptx_file
            for pptx_file in pptx_files
        }

        # Process completed tasks as they finish
        for future in as_completed(future_to_file):
            pptx_file = future_to_file[future]
            try:
                if future.result():
                    successful += 1
                else:
                    failed += 1
            except Exception as e:
                logger.error(f"Unexpected error processing {pptx_file.name}: {str(e)}")
                failed += 1

    # Log summary
    elapsed_time = time.time() - start_time
    logger.info(f"\nProcessing complete:")
    logger.info(f"Time taken: {elapsed_time:.2f} seconds")
    logger.info(f"Successfully processed: {successful} files")
    logger.info(f"Failed to process: {failed} files")

if __name__ == "__main__":
    import sys

    if len(sys.argv) != 2:
        print("Usage: python script.py <directory_path>")
        sys.exit(1)

    directory_path = sys.argv[1]
    process_directory(directory_path)
